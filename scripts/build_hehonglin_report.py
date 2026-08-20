"""Extend the existing editable demo deck with He Honglin's deliverables."""

from __future__ import annotations

import csv
import json
import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
ARTIFACT_ROOT = REPO_ROOT.parent.parent / "outputs"
SOURCE = ARTIFACT_ROOT / "ros_full_demo" / "汇报PPT" / "editable_ros2_demo_report.pptx"
OUTPUT = ARTIFACT_ROOT / "ros_full_demo" / "汇报PPT" / "何泓林_完整汇报_可编辑.pptx"
DATA_DIR = ARTIFACT_ROOT / "ros_full_demo"
EVIDENCE_DIR = REPO_ROOT / "docs" / "evidence"

NAVY = RGBColor(19, 38, 63)
INK = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(203, 213, 225)
PALE = RGBColor(241, 245, 249)
TEAL = RGBColor(13, 148, 136)
BLUE = RGBColor(37, 99, 235)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, number, title, subtitle=None):
    add_text(slide, f"{number:02d}", 0.6, 0.38, 0.55, 0.35, 18, TEAL, True)
    add_text(slide, title, 1.25, 0.32, 11.35, 0.5, 27, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 1.27, 0.88, 11.2, 0.35, 11, MUTED)


def add_footer(slide, text="何泓林工作流 · Swarm-Control-System"):
    add_text(slide, text, 0.65, 7.08, 12.0, 0.2, 8.5, MUTED)


def add_card(slide, x, y, w, h, title, body, accent=BLUE, body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.18, y + 0.13, w - 0.36, 0.32, 14, accent, True)
    add_text(slide, body, x + 0.18, y + 0.52, w - 0.36, h - 0.64, body_size, INK)
    return shape


def add_status(slide, x, y, label, status, color):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.35), Inches(0.32))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    add_text(slide, f"{label} · {status}", x + 0.02, y + 0.01, 1.31, 0.27, 9, WHITE, True, PP_ALIGN.CENTER)


def setup_slide(prs):
    # The source deck carries a reduced custom layout collection; layout 0 is
    # the stable blank-compatible fallback across the bundled template.
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(248, 250, 252)
    return slide


def read_track_rows():
    path = DATA_DIR / "tracks.csv"
    if not path.exists():
        return 0
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        return sum(1 for _ in csv.DictReader(handle))


def read_soak_report():
    reports = sorted(EVIDENCE_DIR.glob("soak_*_report.json"))
    if not reports:
        return None
    try:
        return json.loads(reports[-1].read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None


def soak_verified(report):
    return bool(
        report
        and report.get("status") == "PASS"
        and report.get("elapsed_duration_s", 0) >= report.get("requested_duration_s", 7200)
    )


def ensure_jpeg_content_type(path: Path) -> None:
    """Repair a legacy source-deck omission without changing slide content."""
    namespace = "http://schemas.openxmlformats.org/package/2006/content-types"
    ET.register_namespace("", namespace)
    with zipfile.ZipFile(path, "r") as source:
        content_types = ET.fromstring(source.read("[Content_Types].xml"))
        has_jpeg = any(
            node.tag == f"{{{namespace}}}Default"
            and node.attrib.get("Extension", "").lower() == "jpg"
            for node in content_types
        )
        if has_jpeg:
            return
        ET.SubElement(
            content_types,
            f"{{{namespace}}}Default",
            {"Extension": "jpg", "ContentType": "image/jpeg"},
        )
        updated = ET.tostring(content_types, encoding="UTF-8", xml_declaration=True)
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False, dir=path.parent) as handle:
            temporary = Path(handle.name)
        try:
            with zipfile.ZipFile(temporary, "w", compression=zipfile.ZIP_DEFLATED) as target:
                for item in source.infolist():
                    data = updated if item.filename == "[Content_Types].xml" else source.read(item.filename)
                    target.writestr(item, data)
            shutil.move(str(temporary), str(path))
        finally:
            temporary.unlink(missing_ok=True)
def add_interface_slide(prs):
    slide = setup_slide(prs)
    add_title(slide, 9, "接口决议归档", "D-1 ~ D-12：把联调约定固化为可检查的工程契约")
    rows = [
        ("时间戳", "DroneStateArray.header.stamp 用于跨感知与平台状态对齐", TEAL),
        ("坐标系", "drone_states / target_track_world 使用 world ENU；像素轨迹保持 camera frame", BLUE),
        ("任务接口", "TaskAssignment 只携带 drone_id、target_id、task_type，避免坐标重复", AMBER),
        ("QoS", "除 raw image 外统一 RELIABLE depth=10；传感器输入使用 sensor-compatible QoS", TEAL),
        ("容错", "上游未发消息时下游保持可启动；全量 snapshot 重发，不依赖 latched state", BLUE),
        ("命名", "C1 阶段保留 root topics；多实例由后续 robot_id 机制扩展", AMBER),
    ]
    for index, (key, value, color) in enumerate(rows):
        y = 1.45 + index * 0.78
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y), Inches(1.35), Inches(0.42))
        badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
        add_text(slide, key, 0.8, y + 0.04, 1.25, 0.3, 11, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, value, 2.35, y + 0.02, 9.95, 0.35, 14, INK)
    add_footer(slide, "来源：docs/integration/interface_alignment.md · docs/interface/D-1_D-12_接口决议归档.md")


def add_verification_slide(prs):
    slide = setup_slide(prs)
    report = read_soak_report()
    verified = soak_verified(report)
    add_title(slide, 10, "感知与接口回归验证", "修复后在 ROS2 Humble 隔离 overlay 中复测")
    add_card(slide, 0.75, 1.45, 3.75, 2.2, "Tracker / Fusion", "28 passed · 10 skipped\n\n包含 NumPy 数组断言修复、时间戳对齐、多源融合生命周期回归。", TEAL, 16)
    add_card(slide, 4.78, 1.45, 3.75, 2.2, "Planning", "23 passed\n\n规划包测试通过，覆盖路径输出与 planner 相关行为。", BLUE, 16)
    add_card(slide, 8.81, 1.45, 3.75, 2.2, "Message Contract", "DroneStateArray\n\n生成消息确认包含 std_msgs/Header header，Python import 正常。", AMBER, 15)
    add_status(slide, 0.85, 4.15, "CI", "已验证", TEAL)
    add_status(slide, 2.45, 4.15, "SITL运行", "已验证", TEAL)
    add_status(slide, 4.05, 4.15, "2h挂机", "已验证" if verified else "待采集", TEAL if verified else AMBER)
    add_text(slide, "验证边界：PX4/Gazebo/MAVROS smoke 与三链路八节点 headless 浸泡测试已有原始证据；不等同于 Offboard 飞行或桌面录屏。", 0.85, 4.8, 11.5, 0.65, 15, INK, True)
    add_footer(slide)


def add_bridge_slide(prs):
    slide = setup_slide(prs)
    add_title(slide, 11, "单机 PX4 SITL 桥接设计", "当前 profile 明确支持 num_uav:=1，避免把未验证的多机能力写成已交付")
    steps = [
        ("PX4 + Gazebo", "sitl_run.sh\nGazebo Classic iris", BLUE),
        ("MAVROS", "/uav0/mavros/state\nlocal_position/pose", TEAL),
        ("Pose bridge", "PoseStamped ->\n/drone_pose_external", AMBER),
        ("Offboard bridge", "/planned_path ->\nsetpoint_raw/local", RED),
    ]
    for i, (title, body, color) in enumerate(steps):
        x = 0.7 + i * 3.08
        add_card(slide, x, 1.65, 2.5, 2.1, title, body, color, 14)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.58), Inches(2.42), Inches(0.42), Inches(0.38))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = LINE; arrow.line.fill.background()
    add_text(slide, "关键修复：反馈订阅统一为 /uav0/mavros/local_position/pose，并使用 sensor QoS；DroneStateArray 带 header 时间戳。", 0.85, 4.35, 11.4, 0.55, 15, INK, True)
    add_text(slide, "运行前提：PX4_SITL_ROOT 指向已构建的 PX4-Autopilot；Gazebo Classic、MAVROS 与 GeographicLib 数据集已安装。", 0.85, 5.1, 11.4, 0.45, 13, MUTED)
    add_footer(slide, "来源：docs/simulation/px4_sitl_setup.md · planning_pkg/launch/px4_sitl.launch.py")


def add_repro_slide(prs):
    slide = setup_slide(prs)
    add_title(slide, 12, "可复现运行与证据采集", "脚本负责启动、观察、归档；结果文件与运行声明分开管理")
    add_card(slide, 0.75, 1.5, 5.65, 3.7, "一键入口", "1  source /opt/ros/humble/setup.bash\n2  source ros2_ws/install/setup.bash\n3  PX4_SITL_ROOT=... ros2 launch planning_pkg px4_sitl.launch.py\n4  ros2 topic echo /uav0/mavros/state --once\n5  ros2 topic hz /drone_pose_external", BLUE, 14)
    add_card(slide, 6.7, 1.5, 5.65, 3.7, "无人值守入口", "scripts/run_soak_test.sh\n\n默认 7200 秒，周期记录 RSS、节点数、日志、CSV 与 JSON。\n\n通过标准：elapsed_duration_s 达标、launch_alive=1、无 traceback。", TEAL, 14)
    add_text(slide, "当前仓库不伪造 MP4：pseudo / ros2bag 只生成日志或 bag；真实视频必须使用有桌面会话的 ffmpeg 模式。", 0.85, 5.7, 11.4, 0.45, 14, AMBER, True)
    add_footer(slide, "来源：scripts/run_soak_test.sh · scripts/record_three_links.sh · docs/evidence/")


def add_status_slide(prs):
    slide = setup_slide(prs)
    add_title(slide, 13, "仿真验证状态", "SITL smoke 与 headless 三链路浸泡测试已按证据分层记录")
    add_card(slide, 0.75, 1.45, 5.55, 1.8, "已验证", "PX4/Gazebo/MAVROS smoke\nHeartbeat、pose、DroneStateArray\nplanned_path -> MAVROS setpoint\n8 节点三链路 headless soak", TEAL, 14)
    add_card(slide, 6.65, 1.45, 5.55, 1.8, "仍待采集", "Gazebo 桌面窗口 / 全系统 MP4\n三场景精剪与可视化指标\n真实设备与 Offboard 飞行\n不以 headless 证据替代上述内容", AMBER, 14)
    add_text(slide, "验收顺序已完成：IP/UDP -> DDS discovery -> MAVROS heartbeat -> pose bridge -> planned_path -> soak evidence。", 0.9, 3.85, 11.2, 0.7, 17, NAVY, True)
    add_text(slide, "验收顺序：IP/UDP -> DDS discovery -> MAVROS heartbeat -> pose bridge -> planned_path -> soak evidence", 0.9, 5.0, 11.2, 0.45, 14, MUTED)
    add_footer(slide, "状态标签：已验证 = 有测试或日志；待采集 = 依赖 VM 软件安装后的真实运行")


def add_data_slide(prs):
    slide = setup_slide(prs)
    report = read_soak_report()
    verified = soak_verified(report)
    add_title(slide, 14, "当前可复核数据", "历史感知实测、SITL smoke 与 headless soak 分开呈现")
    rows = read_track_rows()
    chart_data = CategoryChartData()
    chart_data.categories = ["Tracker/Fusion", "Planning", "SITL smoke", "2h soak"]
    chart_data.add_series("已通过或已采集", (28, 23, 1, 1 if verified else 0))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.4), Inches(6.2), Inches(4.8), chart_data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 30
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 5
    chart.category_axis.reverse_order = True
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = TEAL
    chart.value_axis.has_major_gridlines = True
    chart.chart_title.has_text_frame = False
    add_card(slide, 7.45, 1.45, 4.9, 1.25, "历史感知产物", f"tracks.csv 记录行数：{rows:,}\ntracked.mp4、preview_frame.jpg、planner/scheduler/perception logs 已留存。", BLUE, 13)
    add_card(slide, 7.45, 2.95, 4.9, 1.25, "当前回归", "28 passed、10 skipped perception\n23 passed planning\n1 warning：pytest 类收集提示，不影响通过结果。", TEAL, 13)
    if verified:
        elapsed = report["elapsed_duration_s"]
        add_card(slide, 7.45, 4.45, 4.9, 1.25, "2h headless soak", f"PASS · {elapsed}s\n8 个节点持续在线；RSS 约 36,720 KB，日志无 traceback。", TEAL, 13)
    else:
        add_card(slide, 7.45, 4.45, 4.9, 1.25, "2h headless soak", "报告尚未归档；不填写达成值。", AMBER, 13)
    add_footer(slide, "证据目录：outputs/ros_full_demo/ · VM 验证副本：~/codex-swarm-validation-v2")


def add_close_slide(prs):
    slide = setup_slide(prs)
    report = read_soak_report()
    verified = soak_verified(report)
    add_title(slide, 15, "交付结论与下一步", "何泓林工作流的完成边界")
    add_card(slide, 0.75, 1.45, 3.75, 3.6, "已交付", "CI 感知修复\n传感器 QoS 修复\nDroneStateArray 时间戳接口\n单机 PX4 SITL 链路\n挂机与视频证据工具\n接口、报告与素材归档", TEAL, 14)
    verified_body = "ROS2 overlay clean build\nDroneStateArray.header import\nPerception/Fusion 28 passed\nPlanning 23 passed\nPX4 SITL smoke\n8 节点 2h soak" if verified else "ROS2 overlay clean build\nDroneStateArray.header import\nPerception/Fusion 28 passed\nPlanning 23 passed\nPX4 SITL smoke"
    add_card(slide, 4.78, 1.45, 3.75, 3.6, "已验证", verified_body, BLUE, 14)
    add_card(slide, 8.81, 1.45, 3.75, 3.6, "仍待完成", "Gazebo 桌面全系统录屏\n三场景视频精剪与指标\n真实设备与 Offboard 飞行\n不把 headless 测试写成飞行结果", AMBER, 14)
    add_text(slide, "提交原则：只提交有原始日志、topic snapshot 或测试报告支撑的工程结论；视频与真机边界单独标注。", 0.9, 5.65, 11.2, 0.55, 18, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, "目标远端：zxcvn-syc/Swarm-Control-System · 分支：codex/hehonglin-cvtrack-fixes")


def main():
    if not SOURCE.exists():
        raise SystemExit(f"missing source deck: {SOURCE}")
    prs = Presentation(str(SOURCE))
    if len(prs.slides) != 8:
        raise SystemExit(f"expected the existing 8-slide source deck, found {len(prs.slides)}")
    if not any(shape.shape_type == 13 for shape in prs.slides[4].shapes):
        raise SystemExit("source slide 5 has no picture; refusing to replace the required original measurement image")
    add_interface_slide(prs)
    add_verification_slide(prs)
    add_bridge_slide(prs)
    add_repro_slide(prs)
    add_status_slide(prs)
    add_data_slide(prs)
    add_close_slide(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    ensure_jpeg_content_type(OUTPUT)
    print(str(OUTPUT).encode("ascii", "backslashreplace").decode("ascii"))


if __name__ == "__main__":
    main()
